import os
import re
from docx import Document
from docx.shared import Pt, Inches, RGBColor, Cm
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.enum.table import WD_TABLE_ALIGNMENT
from docx.oxml.ns import qn
from docx.oxml import OxmlElement
from docx.text.paragraph import Paragraph

TEMPLATE_STYLES = {
    "corporativo": {
        "font_heading": "Calibri",
        "font_body": "Calibri",
        "heading_color": RGBColor(0, 0, 0),
        "body_color": RGBColor(64, 64, 64),
        "heading_size": 18,
        "body_size": 10,
        "table_header_bg": "2D2D2D",
        "table_header_color": RGBColor(255, 255, 255),
        "table_row_alt": "F2F2F2",
        "accent_color": RGBColor(250, 128, 114),
    },
    "academico": {
        "font_heading": "Georgia",
        "font_body": "Garamond",
        "heading_color": RGBColor(26, 26, 46),
        "body_color": RGBColor(55, 65, 81),
        "heading_size": 18,
        "body_size": 11,
        "table_header_bg": "1A1A2E",
        "table_header_color": RGBColor(255, 255, 255),
        "table_row_alt": "EEF0F5",
        "accent_color": RGBColor(30, 64, 175),
    },
    "minimalista": {
        "font_heading": "Arial",
        "font_body": "Arial",
        "heading_color": RGBColor(17, 24, 39),
        "body_color": RGBColor(75, 85, 99),
        "heading_size": 16,
        "body_size": 10,
        "table_header_bg": "F3F4F6",
        "table_header_color": RGBColor(17, 24, 39),
        "table_row_alt": "FAFAFA",
        "accent_color": RGBColor(107, 114, 128),
    },
}

def _get_style_config(template_style):
    return TEMPLATE_STYLES.get(template_style, TEMPLATE_STYLES["corporativo"])

def _add_free_watermark_header(doc):
    watermark_text = "Generado por Pertinentia\u00AE \u2014 pertinentia.com"
    for section in doc.sections:
        for hdr_attr in ['header', 'first_page_header', 'even_page_header']:
            try:
                header = getattr(section, hdr_attr)
                header.is_linked_to_previous = False
                p = header.add_paragraph()
                p.alignment = WD_ALIGN_PARAGRAPH.RIGHT
                run = p.add_run(watermark_text)
                run.font.size = Pt(9)
                run.font.color.rgb = RGBColor(192, 192, 192)
                run.font.name = "Calibri"
                run.italic = True
            except Exception:
                pass

DOCS_DIR = "plantillas"
GENERATED_DIR = "generated_docs"
TEMPLATES_MAP = {
    "diagnostica": "EVALUACION_DIAGNOSTICA_301.docx",
    "sumativa": "EVALUACION_SUMATIVA_301.docx",
    "reaccion": "EVALUACION_REACCION_301.docx",
    "guia_observacion": "INSTRUMENTO_GUÍA_DE_OBSERVACIÓN_301.docx",
    "lista_cotejo": "INSTRUMENTO_LISTA_DE_COTEJO_301.docx",
    "hojas_respuestas": "HOJAS_DE_RESPUESTAS.docx",
    "lista_requerimientos": "LISTA_DE_REQUERIMIENTOS_301.docx",
    "contrato_aprendizaje": "CONTRATO_DE_APRENDIZAJE_301.docx",
}

FORBIDDEN_TERMS = [
    (r'(?i)curtido\s+de\s+piel\s+de\s+conejo\s+con\s+pelo', '[Tema del curso]'),
    (r'(?i)curtido\s+de\s+piel\s+de\s+conejo', '[Tema del curso]'),
    (r'(?i)curtido\s+de\s+piel', '[Tema del curso]'),
    (r'(?i)piel\s+de\s+conejo', '[producto del curso]'),
    (r'(?i)Braulio\s+G[oó]mez', '[Nombre del instructor]'),
    (r'(?i)Granja\s+Santez', '[Lugar de imparticion]'),
]

def _sanitize_text(text):
    for pattern, replacement in FORBIDDEN_TERMS:
        text = re.sub(pattern, replacement, text)
    text = re.sub(r'\*+', '', text)
    return text

# PALETA MARCA BLANCA (WHITE-LABEL)
NEUTRAL_BLACK = RGBColor(0, 0, 0)
NEUTRAL_DARK_GRAY = RGBColor(64, 64, 64)
HEADER_BG = "D9D9D9" # Gris estándar profesional para tablas
ROW_ALT_BG = "F2F2F2" # Gris muy claro para filas alternas
WHITE = RGBColor(255, 255, 255)
BLACK = RGBColor(0, 0, 0)

def _set_cell_shading(cell, color_hex):
    shading = cell._element.get_or_add_tcPr()
    shading_elm = shading.makeelement(qn('w:shd'), {
        qn('w:fill'): color_hex,
        qn('w:val'): 'clear'
    })
    shading.append(shading_elm)

def _set_cell_text(cell, text, bold=False, size=10, color=None, align=None, font_name='Calibri'):
    cell.text = ""
    p = cell.paragraphs[0]
    if align:
        p.alignment = align
    # Convierte <br>, <br/>, <br /> (que la IA genera como marcador de salto de
    # línea) en saltos de línea reales dentro de la celda, en vez de insertarlos
    # como texto literal visible.
    normalized = re.sub(r'<br\s*/?>', '\n', str(text), flags=re.IGNORECASE)
    lines = normalized.split('\n')
    for i, line in enumerate(lines):
        run = p.add_run(line)
        run.bold = bold
        run.font.size = Pt(size)
        run.font.name = font_name
        if color:
            run.font.color.rgb = color
        if i < len(lines) - 1:
            run.add_break()

def _add_professional_table(doc, headers, rows):
    table = doc.add_table(rows=1 + len(rows), cols=len(headers))
    table.alignment = WD_TABLE_ALIGNMENT.CENTER
    table.style = 'Table Grid'

    header_row = table.rows[0]
    for i, header_text in enumerate(headers):
        cell = header_row.cells[i]
        _set_cell_text(cell, header_text, bold=True, size=9, color=BLACK, align=WD_ALIGN_PARAGRAPH.CENTER)
        _set_cell_shading(cell, HEADER_BG)

    for row_idx, row_data in enumerate(rows):
        row = table.rows[row_idx + 1]
        bg = ROW_ALT_BG if row_idx % 2 == 0 else "FFFFFF"
        for col_idx, cell_text in enumerate(row_data):
            cell = row.cells[col_idx]
            _set_cell_text(cell, str(cell_text), size=9, color=NEUTRAL_BLACK)
            _set_cell_shading(cell, bg)

    return table

def _detect_markdown_table(lines, start_idx):
    if start_idx >= len(lines):
        return None, start_idx
    line = lines[start_idx].strip()
    if '|' not in line:
        return None, start_idx

    headers = [c.strip() for c in line.strip('|').split('|')]
    if not headers or len(headers) < 2:
        return None, start_idx

    next_idx = start_idx + 1
    if next_idx < len(lines) and re.match(r'^[\s|:-]+$', lines[next_idx].strip()):
        next_idx += 1

    rows = []
    while next_idx < len(lines):
        row_line = lines[next_idx].strip()
        if '|' not in row_line or row_line.startswith('#'):
            break
        cells = [c.strip() for c in row_line.strip('|').split('|')]
        if len(cells) >= 2:
            while len(cells) < len(headers):
                cells.append("")
            rows.append(cells[:len(headers)])
        next_idx += 1

    if rows:
        return {"headers": headers, "rows": rows}, next_idx
    return None, start_idx

def _add_logo_to_header(doc, logo_path):
    """Inserta el logotipo de la organizacion en el encabezado de cada pagina.
    Funciona para todos los entregables sin alterar el cuerpo del documento.
    Devuelve True si se inserto al menos una vez; si no hay logo o falla, no
    rompe la generacion (degradacion silenciosa, sin placeholder en encabezado)."""
    if not logo_path:
        return False
    try:
        if not os.path.exists(logo_path):
            return False
    except Exception:
        return False
    inserted = False
    for section in doc.sections:
        for hdr_attr in ['header', 'first_page_header', 'even_page_header']:
            try:
                header = getattr(section, hdr_attr)
                header.is_linked_to_previous = False
                p = header.add_paragraph()
                p.alignment = WD_ALIGN_PARAGRAPH.LEFT
                p.add_run().add_picture(logo_path, height=Inches(0.45))
                inserted = True
            except Exception:
                pass
    return inserted

def generate_from_template(template_key, content_data, user_tier="FREE", logo_path=None):
    os.makedirs(GENERATED_DIR, exist_ok=True)
    template_name = TEMPLATES_MAP.get(template_key)
    if not template_name:
        return None, f"Plantilla '{template_key}' no encontrada"

    template_path = os.path.join(DOCS_DIR, template_name)
    if not os.path.exists(template_path):
        return None, f"Archivo de plantilla no existe: {template_path}"

    try:
        doc = Document(template_path)
        course_info = content_data.get("course_info", {})
        _fill_header_fields(doc, course_info)
        if template_key in ("diagnostica", "sumativa"):
            _fill_questionnaire(doc, content_data)
        elif template_key in ("guia_observacion", "lista_cotejo"):
            _fill_observation_guide(doc, content_data)
        elif template_key == "hojas_respuestas":
            _fill_answer_sheet(doc, content_data)
        elif template_key == "contrato_aprendizaje":
            _fill_contrato(doc, content_data)
        elif template_key == "lista_requerimientos":
            _fill_lista_requerimientos(doc, content_data)

        _add_logo_to_header(doc, logo_path)

        if user_tier == "FREE":
            _add_free_watermark_header(doc)

        curso_safe = re.sub(r'[^\w\s-]', '', course_info.get("nombre_curso", "curso")).strip().replace(' ', '_')[:30]
        output_name = f"{template_key}_{curso_safe}.docx"
        output_path = os.path.join(GENERATED_DIR, output_name)
        doc.save(output_path)
        return output_path, None
    except Exception as e:
        return None, str(e)

def _fill_header_fields(doc, course_info):
    """Rellena la tabla de datos del curso escribiendo cada VALOR en la celda
    BLANCA contigua a su etiqueta (celda de color), nunca dentro de la etiqueta.

    Las plantillas comparten un encabezado de 4 filas × 8 columnas: las etiquetas
    viven en celdas de color y el valor pertenece a la celda contigua a la derecha.
    """
    if not course_info:
        return

    def _norm(s):
        return re.sub(r'\s+', ' ', (s or '').strip().lower()).rstrip(':').strip()

    label_values = {
        _norm("Nombre del Curso-Taller"): course_info.get("nombre_curso", ""),
        _norm("Nombre del diseñador"): course_info.get("nombre_disenador", ""),
        _norm("Nombre del instructor"): course_info.get("nombre_instructor", ""),
        _norm("Lugar de Impartición"): course_info.get("lugar", ""),
        _norm("Duración del curso"): course_info.get("duracion", ""),
        _norm("Horario"): course_info.get("horario", ""),
        _norm("Fecha de impartición"): course_info.get("fecha", ""),
    }

    for table in doc.tables:
        for row in table.rows:
            uniq = []
            seen = set()
            for c in row.cells:
                cid = id(c._tc)
                if cid in seen:
                    continue
                seen.add(cid)
                uniq.append(c)
            for i, cell in enumerate(uniq):
                value = label_values.get(_norm(cell.text))
                if not value or str(value).strip().startswith('['):
                    continue
                for j in range(i + 1, len(uniq)):
                    if not uniq[j].text.strip():
                        _set_cell_text(uniq[j], str(value), size=10, color=NEUTRAL_BLACK)
                        break

def _fill_questionnaire(doc, content_data):
    reactivos = content_data.get("reactivos", [])
    instrucciones_instructor = content_data.get("instrucciones_instructor", "")
    instrucciones_participante = content_data.get("instrucciones_participante", "")

    for paragraph in doc.paragraphs:
        if "Redactar las instrucciones para el instructor" in paragraph.text:
            if instrucciones_instructor:
                paragraph.clear()
                run = paragraph.add_run(f"Indicaciones para el instructor: {instrucciones_instructor}")
                run.font.size = Pt(10)
        elif "Redactar las instrucciones para resolver" in paragraph.text:
            if instrucciones_participante:
                paragraph.clear()
                run = paragraph.add_run(f"Indicaciones para el participante: {instrucciones_participante}")
                run.font.size = Pt(10)
        elif "Diseñar el cuestionario" in paragraph.text:
            paragraph.clear()
            for i, reactivo in enumerate(reactivos, 1):
                pregunta = reactivo.get("pregunta", "")
                valor = reactivo.get("valor", "")
                run = paragraph.add_run(f"\n{i}. {pregunta} (Valor: {valor})\n")
                run.font.size = Pt(10)

def _fill_observation_guide(doc, content_data):
    reactivos = content_data.get("reactivos", [])
    for table in doc.tables:
        for i, row in enumerate(table.rows):
            if i == 0:
                continue
            row_idx = i - 1
            if row_idx < len(reactivos):
                reactivo = reactivos[row_idx]
                if len(row.cells) >= 2:
                    desc_cell = row.cells[1]
                    if "Redactar" in desc_cell.text:
                        _set_cell_text(desc_cell, reactivo.get("descripcion", ""))
                if len(row.cells) >= 4:
                    valor_cell = row.cells[3] if len(row.cells) > 3 else row.cells[-1]
                    _set_cell_text(valor_cell, str(reactivo.get("valor", "")))

def _fill_answer_sheet(doc, content_data):
    respuestas = content_data.get("respuestas", [])
    for paragraph in doc.paragraphs:
        if "Redactar las respuestas" in paragraph.text:
            paragraph.clear()
            for resp in respuestas:
                evaluacion = resp.get("evaluacion", "")
                answers = resp.get("respuestas", [])
                run = paragraph.add_run(f"\n{evaluacion}\n")
                run.bold = True
                run.font.size = Pt(11)
                for j, ans in enumerate(answers, 1):
                    run2 = paragraph.add_run(f"{j}. {ans}\n")
                    run2.font.size = Pt(10)
            break

def _fill_contrato(doc, content_data):
    compromisos_instructor = content_data.get("compromisos_instructor", [])
    compromisos_participante = content_data.get("compromisos_participante", [])
    contenido_completo = content_data.get("contenido_completo", "")
    course_info = content_data.get("course_info", {})
    objetivo_curso = content_data.get("objetivo_curso", "")
    criterios_evaluacion = content_data.get("criterios_evaluacion", "")
    derechos_participante = content_data.get("derechos_participante", [])

    # El encabezado (datos del curso) lo rellena _fill_header_fields en la celda
    # blanca correcta; aquí solo se llenan compromisos, objetivo y cláusulas.

    instructor_table = None
    participant_table = None
    for i, paragraph in enumerate(doc.paragraphs):
        p_lower = paragraph.text.lower().strip()
        if "compromiso" in p_lower and "instructor" in p_lower:
            for ti, table in enumerate(doc.tables):
                tbl_elem = table._tbl
                para_elem = paragraph._p
                if para_elem.getparent() == tbl_elem.getparent():
                    parent = para_elem.getparent()
                    para_idx = list(parent).index(para_elem)
                    for sibling in list(parent)[para_idx+1:]:
                        if sibling.tag.endswith('}tbl'):
                            for t in doc.tables:
                                if t._tbl is sibling:
                                    instructor_table = t
                                    break
                            break
                        if sibling.tag.endswith('}p'):
                            sib_text = sibling.text or ""
                            if "participante" in sib_text.lower():
                                break
                    break
        if "compromiso" in p_lower and "participante" in p_lower:
            for ti, table in enumerate(doc.tables):
                tbl_elem = table._tbl
                para_elem = paragraph._p
                parent = para_elem.getparent()
                para_idx = list(parent).index(para_elem)
                for sibling in list(parent)[para_idx+1:]:
                    if sibling.tag.endswith('}tbl'):
                        for t in doc.tables:
                            if t._tbl is sibling:
                                participant_table = t
                                break
                        break
                break

    if instructor_table is None and len(doc.tables) >= 4:
        instructor_table = doc.tables[2]
    if participant_table is None and len(doc.tables) >= 4:
        participant_table = doc.tables[3]

    def _fill_commitment_table(table, commitments):
        if not table or not commitments:
            return
        if len(table.rows) < 2:
            return
        num_cell = table.rows[1].cells[0]
        desc_cell = table.rows[1].cells[1] if len(table.rows[1].cells) >= 2 else None
        for p in num_cell.paragraphs:
            p.clear()
        first_para = num_cell.paragraphs[0]
        for ci, c in enumerate(commitments[:6]):
            r = first_para.add_run(str(ci + 1))
            r.font.size = Pt(9)
            r.font.color.rgb = NEUTRAL_DARK_GRAY
            if ci < len(commitments[:6]) - 1:
                first_para.add_run("\n")
        if desc_cell:
            for p in desc_cell.paragraphs:
                p.clear()
            first_para_desc = desc_cell.paragraphs[0]
            for ci, c in enumerate(commitments[:6]):
                text = _sanitize_text(c)
                r = first_para_desc.add_run(text)
                r.font.size = Pt(9)
                r.font.color.rgb = NEUTRAL_DARK_GRAY
                if ci < len(commitments[:6]) - 1:
                    first_para_desc.add_run("\n")

    _fill_commitment_table(instructor_table, compromisos_instructor)
    _fill_commitment_table(participant_table, compromisos_participante)

    insert_before_signatures = None
    for i, paragraph in enumerate(doc.paragraphs):
        if "____" in paragraph.text or "Firma" in paragraph.text:
            insert_before_signatures = paragraph
            break

    if objetivo_curso:
        p = OxmlElement('w:p')
        if insert_before_signatures:
            insert_before_signatures._p.getparent().insert(
                list(insert_before_signatures._p.getparent()).index(insert_before_signatures._p), p
            )
        else:
            doc.element.body.append(p)
        new_para = Paragraph(p, doc)
        r_title = new_para.add_run("Objetivo General del Curso:\n")
        r_title.bold = True
        r_title.font.size = Pt(11)
        r_title.font.color.rgb = NEUTRAL_BLACK
        r_content = new_para.add_run(_sanitize_text(objetivo_curso))
        r_content.font.size = Pt(10)
        r_content.font.color.rgb = NEUTRAL_DARK_GRAY

    if criterios_evaluacion:
        p = OxmlElement('w:p')
        if insert_before_signatures:
            insert_before_signatures._p.getparent().insert(
                list(insert_before_signatures._p.getparent()).index(insert_before_signatures._p), p
            )
        else:
            doc.element.body.append(p)
        new_para = Paragraph(p, doc)
        r_title = new_para.add_run("\nCriterios y Momentos de Evaluación:\n")
        r_title.bold = True
        r_title.font.size = Pt(11)
        r_title.font.color.rgb = NEUTRAL_BLACK
        r_content = new_para.add_run(_sanitize_text(criterios_evaluacion))
        r_content.font.size = Pt(10)
        r_content.font.color.rgb = NEUTRAL_DARK_GRAY

    if derechos_participante:
        p = OxmlElement('w:p')
        if insert_before_signatures:
            insert_before_signatures._p.getparent().insert(
                list(insert_before_signatures._p.getparent()).index(insert_before_signatures._p), p
            )
        else:
            doc.element.body.append(p)
        new_para = Paragraph(p, doc)
        r_title = new_para.add_run("\nDerechos del Participante:\n")
        r_title.bold = True
        r_title.font.size = Pt(11)
        r_title.font.color.rgb = NEUTRAL_BLACK
        for di, d in enumerate(derechos_participante):
            r_item = new_para.add_run(f"{di+1}. {_sanitize_text(d)}\n")
            r_item.font.size = Pt(10)
            r_item.font.color.rgb = NEUTRAL_DARK_GRAY

    has_structured = bool(compromisos_instructor or compromisos_participante or objetivo_curso)
    if not has_structured and contenido_completo and len(contenido_completo) > 50:
        p = OxmlElement('w:p')
        if insert_before_signatures:
            insert_before_signatures._p.getparent().insert(
                list(insert_before_signatures._p.getparent()).index(insert_before_signatures._p), p
            )
        else:
            doc.element.body.append(p)
        new_para = Paragraph(p, doc)
        for line in contenido_completo.split("\n"):
            line = _sanitize_text(line.strip())
            if not line:
                continue
            r = new_para.add_run(f"{line}\n")
            r.font.size = Pt(10)
            r.font.color.rgb = NEUTRAL_DARK_GRAY

_CALLOUT_PATTERNS = {
    'nota': ('\U0001F4DD', 'NOTA', 'D6E4F0'),
    'importante': ('\u26A0\uFE0F', 'IMPORTANTE', 'FCE4D6'),
    'tip': ('\U0001F4A1', 'RECOMENDACI\u00D3N', 'E2EFDA'),
    'recomendacion': ('\U0001F4A1', 'RECOMENDACI\u00D3N', 'E2EFDA'),
    'ejemplo': ('\U0001F4CC', 'EJEMPLO', 'E8D5F5'),
    'advertencia': ('\u26A0\uFE0F', 'ADVERTENCIA', 'FCE4D6'),
}

def _detect_callout(line):
    import unicodedata
    stripped = line.strip()
    if not stripped:
        return None
    if stripped.startswith('> '):
        return ('nota', stripped[2:].strip())
    bold_match = re.match(r'^\*\*([^*]+)\*\*\s*(.*)', stripped)
    if bold_match:
        bold_inner = bold_match.group(1).strip()
        rest = bold_match.group(2).strip()
        if bold_inner.endswith(':'):
            label = bold_inner[:-1].strip()
            norm_label = ''.join(c for c in unicodedata.normalize('NFKD', label.lower()) if not unicodedata.category(c).startswith('M'))
            for key in _CALLOUT_PATTERNS:
                if norm_label == key:
                    return (key, rest)
    return None

def _add_callout_box(doc, callout_type, content_text, sc):
    info = _CALLOUT_PATTERNS.get(callout_type, _CALLOUT_PATTERNS['nota'])
    icon, label, bg_hex = info
    table = doc.add_table(rows=1, cols=1)
    table.alignment = WD_TABLE_ALIGNMENT.CENTER
    table.style = 'Table Grid'
    cell = table.rows[0].cells[0]
    _set_cell_shading(cell, bg_hex)
    tbl = table._tbl
    tblPr = tbl.tblPr if tbl.tblPr is not None else OxmlElement('w:tblPr')
    tblW = OxmlElement('w:tblW')
    tblW.set(qn('w:w'), '5000')
    tblW.set(qn('w:type'), 'pct')
    tblPr.append(tblW)
    cell.text = ""
    header_p = cell.paragraphs[0]
    header_run = header_p.add_run(f"{icon} {label}")
    header_run.bold = True
    header_run.font.size = Pt(9)
    header_run.font.name = sc["font_heading"]
    header_run.font.color.rgb = sc["heading_color"]
    if content_text:
        content_p = cell.add_paragraph()
        content_run = content_p.add_run(content_text)
        content_run.font.size = Pt(sc["body_size"])
        content_run.font.name = sc["font_body"]
        content_run.font.color.rgb = sc["body_color"]
    doc.add_paragraph("").paragraph_format.space_before = Pt(2)

def _add_docx_cover_page(doc, title, course_name, sc, logo_path=None):
    doc.add_paragraph("").paragraph_format.space_before = Pt(60)
    cover_table = doc.add_table(rows=1, cols=1)
    cover_table.alignment = WD_TABLE_ALIGNMENT.CENTER
    cover_cell = cover_table.rows[0].cells[0]
    _set_cell_shading(cover_cell, sc["table_header_bg"])
    tbl = cover_table._tbl
    tblPr = tbl.tblPr if tbl.tblPr is not None else OxmlElement('w:tblPr')
    tblW = OxmlElement('w:tblW')
    tblW.set(qn('w:w'), '5000')
    tblW.set(qn('w:type'), 'pct')
    tblPr.append(tblW)
    cover_cell.text = ""
    spacer_p = cover_cell.paragraphs[0]
    spacer_p.add_run("").font.size = Pt(20)
    display_title = title.replace('_', ' ')
    title_p = cover_cell.add_paragraph()
    title_p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    title_run = title_p.add_run(display_title)
    title_run.font.size = Pt(26)
    title_run.bold = True
    title_run.font.color.rgb = sc["table_header_color"]
    title_run.font.name = sc["font_heading"]
    title_p.paragraph_format.space_before = Pt(24)
    sep_p = cover_cell.add_paragraph()
    sep_p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    sep_run = sep_p.add_run("\u2500" * 40)
    sep_run.font.size = Pt(10)
    sep_run.font.color.rgb = sc.get("accent_color", RGBColor(250, 128, 114))
    if course_name:
        course_p = cover_cell.add_paragraph()
        course_p.alignment = WD_ALIGN_PARAGRAPH.CENTER
        course_run = course_p.add_run(course_name)
        course_run.font.size = Pt(18)
        course_run.bold = True
        course_run.font.color.rgb = sc["table_header_color"]
        course_run.font.name = sc["font_heading"]
        course_p.paragraph_format.space_before = Pt(12)
    subtitle_p = cover_cell.add_paragraph()
    subtitle_p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    subtitle_run = subtitle_p.add_run("Curso-Taller de Capacitaci\u00F3n Presencial Grupal")
    subtitle_run.font.size = Pt(12)
    subtitle_run.italic = True
    subtitle_run.font.color.rgb = sc.get("accent_color", RGBColor(250, 128, 114))
    subtitle_run.font.name = sc["font_body"]
    subtitle_p.paragraph_format.space_before = Pt(8)
    end_spacer = cover_cell.add_paragraph()
    end_spacer.add_run("").font.size = Pt(20)
    doc.add_paragraph("")
    info_fields = [
        ("Instructor:", "[Nombre del Instructor]"),
        ("Dise\u00F1ador:", "[Nombre del Dise\u00F1ador]"),
        ("Fecha:", "[Fecha de impartici\u00F3n]"),
        ("Lugar:", "[Lugar de impartici\u00F3n]"),
        ("Horario:", "[Horario]"),
    ]
    info_table = doc.add_table(rows=len(info_fields), cols=2)
    info_table.alignment = WD_TABLE_ALIGNMENT.CENTER
    for idx, (label, placeholder) in enumerate(info_fields):
        label_cell = info_table.rows[idx].cells[0]
        value_cell = info_table.rows[idx].cells[1]
        label_cell.text = ""
        lp = label_cell.paragraphs[0]
        lr = lp.add_run(label)
        lr.bold = True
        lr.font.size = Pt(10)
        lr.font.name = sc["font_body"]
        lr.font.color.rgb = sc["heading_color"]
        lp.alignment = WD_ALIGN_PARAGRAPH.RIGHT
        value_cell.text = ""
        vp = value_cell.paragraphs[0]
        vr = vp.add_run(placeholder)
        vr.font.size = Pt(10)
        vr.font.name = sc["font_body"]
        vr.font.color.rgb = sc["body_color"]
        vr.italic = True
    for row in info_table.rows:
        for cell in row.cells:
            tc = cell._element
            tcPr = tc.get_or_add_tcPr()
            tcBorders = OxmlElement('w:tcBorders')
            for edge in ['top', 'left', 'bottom', 'right']:
                el = OxmlElement(f'w:{edge}')
                el.set(qn('w:val'), 'none')
                el.set(qn('w:sz'), '0')
                el.set(qn('w:space'), '0')
                tcBorders.append(el)
            tcPr.append(tcBorders)
    doc.add_paragraph("")
    _docx_logo_ok = False
    if logo_path:
        try:
            if os.path.exists(logo_path):
                logo_p = doc.add_paragraph()
                logo_p.alignment = WD_ALIGN_PARAGRAPH.CENTER
                logo_p.add_run().add_picture(logo_path, height=Inches(1.1))
                _docx_logo_ok = True
        except Exception:
            _docx_logo_ok = False
    if not _docx_logo_ok:
        logo_p = doc.add_paragraph()
        logo_p.alignment = WD_ALIGN_PARAGRAPH.CENTER
        logo_run = logo_p.add_run("[ Inserte aqu\u00ED su logo ]")
        logo_run.font.size = Pt(10)
        logo_run.italic = True
        logo_run.font.color.rgb = RGBColor(160, 160, 160)
        logo_run.font.name = sc["font_body"]
    doc.add_page_break()

def _fill_lista_requerimientos(doc, content_data):
    items_by_section = content_data.get("items_by_section", {}) or {}

    # El encabezado (datos del curso) lo rellena _fill_header_fields en la celda
    # blanca correcta; aquí solo se llenan las secciones de requerimientos.

    section_keywords = [
        ("instalaciones", ["instalaciones", "mobiliario"]),
        ("equipo_apoyo", ["equipo de apoyo"]),
        ("materiales_didacticos", ["materiales didácticos", "material didáctico", "materiales didactic"]),
        ("humanos", ["requerimientos humanos"]),
        ("otros", ["otros requerimientos"]),
    ]

    body = doc.element.body
    current_section = None
    for el in list(body.iterchildren()):
        tag = el.tag.split("}")[1] if "}" in el.tag else el.tag
        if tag == "p":
            text = "".join(el.itertext()).strip().lower()
            if not text:
                continue
            for skey, kws in section_keywords:
                if any(kw in text for kw in kws):
                    current_section = skey
                    break
        elif tag == "tbl":
            if current_section is None:
                continue
            target_table = None
            for t in doc.tables:
                if t._tbl is el:
                    target_table = t
                    break
            if target_table is None:
                current_section = None
                continue
            items = items_by_section.get(current_section, []) or []
            data_rows = target_table.rows[1:]
            for i, item in enumerate(items):
                if i >= len(data_rows):
                    break
                if len(data_rows[i].cells) < 2:
                    continue
                desc_cell = data_rows[i].cells[1]
                for p in desc_cell.paragraphs:
                    for r in list(p.runs):
                        r.text = ""
                first_para = desc_cell.paragraphs[0]
                clean_item = _sanitize_text(item)
                if first_para.runs:
                    first_para.runs[0].text = clean_item
                    first_para.runs[0].font.size = Pt(9)
                    first_para.runs[0].font.color.rgb = NEUTRAL_DARK_GRAY
                else:
                    r = first_para.add_run(clean_item)
                    r.font.size = Pt(9)
                    r.font.color.rgb = NEUTRAL_DARK_GRAY
            current_section = None

def generate_custom_docx(title, sections, course_name="", user_tier="FREE", template_style="corporativo", logo_path=None):
    os.makedirs(GENERATED_DIR, exist_ok=True)
    doc = Document()

    is_paid = user_tier in ("PRO", "PREMIUM")
    sc = _get_style_config(template_style) if is_paid else _get_style_config("corporativo")

    style = doc.styles['Normal']
    style.font.name = sc["font_body"]
    style.font.size = Pt(sc["body_size"])
    style.font.color.rgb = sc["body_color"]

    _add_logo_to_header(doc, logo_path)

    is_manual = any(k in title.lower() for k in ['manual', 'instructor', 'participante'])
    if is_manual and course_name:
        _add_docx_cover_page(doc, title, course_name, sc, logo_path=logo_path)
    else:
        heading = doc.add_heading('', level=0)
        heading.alignment = WD_ALIGN_PARAGRAPH.CENTER
        run = heading.add_run(title.replace('_', ' '))
        run.font.color.rgb = sc["heading_color"]
        run.font.size = Pt(sc["heading_size"])
        run.font.name = sc["font_heading"]
        run.bold = True
        if course_name:
            p = doc.add_paragraph()
            p.alignment = WD_ALIGN_PARAGRAPH.CENTER
            run = p.add_run(f"Curso-Taller: {course_name}")
            run.font.size = Pt(12)
            run.font.color.rgb = sc["body_color"]
            run.font.name = sc["font_body"]
            run.italic = True
        doc.add_paragraph("")

    duration_tracker = []
    for section in sections:
        section_title = _sanitize_text(section.get("titulo", ""))
        section_content = _sanitize_text(section.get("contenido", ""))

        if section_title:
            h = doc.add_heading('', level=1)
            run = h.add_run(section_title)
            run.font.color.rgb = sc["heading_color"]
            run.font.size = Pt(14)
            run.font.name = sc["font_heading"]
            run.bold = True

        if section_content:
            content_lines = section_content.split("\n")
            i = 0
            list_number_counter = [0]
            while i < len(content_lines):
                line = content_lines[i].strip()
                if not line:
                    i += 1
                    continue

                table_data, new_idx = _detect_markdown_table(content_lines, i)
                if table_data:
                    _add_styled_table(doc, table_data["headers"], table_data["rows"], sc, duration_tracker=duration_tracker)
                    doc.add_paragraph("")
                    i = new_idx
                    continue

                callout_result = _detect_callout(line)
                if callout_result:
                    callout_type, callout_content = callout_result
                    extra_lines = []
                    j = i + 1
                    while j < len(content_lines):
                        next_line = content_lines[j].strip()
                        if not next_line or next_line.startswith('#') or next_line.startswith('|') or _detect_callout(next_line):
                            break
                        if next_line.startswith('> '):
                            extra_lines.append(next_line[2:].strip())
                            j += 1
                        else:
                            break
                    full_content = (callout_content + ' ' + ' '.join(extra_lines)).strip() if extra_lines else callout_content
                    _add_callout_box(doc, callout_type, full_content, sc)
                    i = j
                    continue

                if line.startswith("### "):
                    list_number_counter[0] = 0
                    h = doc.add_heading('', level=3)
                    run = h.add_run(line[4:])
                    run.font.color.rgb = sc["body_color"]
                    run.font.size = Pt(11)
                    run.font.name = sc["font_heading"]
                elif line.startswith("## "):
                    list_number_counter[0] = 0
                    h = doc.add_heading('', level=2)
                    run = h.add_run(line[3:])
                    run.font.color.rgb = sc["heading_color"]
                    run.font.size = Pt(12)
                    run.font.name = sc["font_heading"]
                    h.paragraph_format.keep_with_next = True
                elif line.startswith("# "):
                    list_number_counter[0] = 0
                    h = doc.add_heading('', level=1)
                    run = h.add_run(line[2:])
                    run.font.color.rgb = sc["heading_color"]
                    run.font.size = Pt(14)
                    run.font.name = sc["font_heading"]
                    h.paragraph_format.keep_with_next = True
                elif line.startswith("- ") or line.startswith("* "):
                    doc.add_paragraph(line[2:], style='List Bullet')
                elif re.match(r'^\d+[\.\)]\s', line):
                    clean = re.sub(r'^\d+[\.\)]\s+', '', line)
                    list_number_counter[0] += 1
                    p = doc.add_paragraph()
                    p.paragraph_format.left_indent = Inches(0.25)
                    p.paragraph_format.first_line_indent = Inches(-0.25)
                    p.add_run(f"{list_number_counter[0]}. {clean}")
                elif line.startswith("**") and line.endswith("**"):
                    p = doc.add_paragraph()
                    run = p.add_run(line.strip("*"))
                    run.bold = True
                    run.font.color.rgb = sc["heading_color"]
                elif line.startswith("| "):
                    pass
                else:
                    doc.add_paragraph(line)
                i += 1

    if duration_tracker and len(duration_tracker) >= 2:
        grand_total = sum(duration_tracker)
        gt_int = int(grand_total) if grand_total == int(grand_total) else round(grand_total, 1)
        hours = gt_int / 60.0
        hours_str = f"{hours:.1f}".rstrip('0').rstrip('.')
        doc.add_paragraph("")
        total_table = doc.add_table(rows=1, cols=1)
        total_table.alignment = WD_TABLE_ALIGNMENT.CENTER
        total_table.style = 'Table Grid'
        tcell = total_table.rows[0].cells[0]
        _set_cell_shading(tcell, sc["table_header_bg"])
        _set_cell_text(tcell, f"TOTAL DE DURACION DEL CURSO: {gt_int} min ({hours_str} h)",
                       bold=True, size=11, color=sc["table_header_color"],
                       align=WD_ALIGN_PARAGRAPH.CENTER, font_name=sc["font_heading"])
        doc.add_paragraph("")

    if user_tier == "FREE":
        _add_free_watermark_header(doc)

    safe_title = re.sub(r'[^\w\s-]', '', title).strip().replace(' ', '_')
    topic_suffix = ""
    if course_name:
        clean_topic = re.sub(r'[^\w\s-]', '', course_name, flags=re.UNICODE).strip().replace(' ', '_')[:40]
        if clean_topic:
            topic_suffix = f"_{clean_topic}"
    output_name = f"{safe_title}{topic_suffix}.docx"
    output_path = os.path.join(GENERATED_DIR, output_name)
    doc.save(output_path)
    return output_path

def _compute_column_widths(headers):
    n = len(headers)
    if n == 0:
        return []
    weights = []
    for h in headers:
        h_low = h.lower().strip()
        if any(k in h_low for k in ['duraci', 'tiempo', 'cantidad', 'ponderaci', '%', 'valor']):
            weights.append(0.7)
        elif any(k in h_low for k in ['tecnica', 'técnica']):
            weights.append(1.3)
        elif any(k in h_low for k in ['momento', 'tipo']):
            weights.append(1.0)
        elif any(k in h_low for k in ['material', 'equipo', 'instrumento']):
            weights.append(1.6)
        elif any(k in h_low for k in ['actividad', 'descripci', 'contenido']):
            weights.append(4.0)
        elif 'tema' in h_low or 'subtema' in h_low:
            weights.append(1.5)
        else:
            weights.append(1.2)
    total_inches = 6.5
    total_w = sum(weights)
    return [Inches(total_inches * w / total_w) for w in weights]

def _extract_minutes(cell_text):
    s = str(cell_text).strip().lower()
    if not s:
        return None
    m = re.search(r'(\d+(?:\.\d+)?)\s*(h(?:r|rs|oras?)?|min(?:utos?)?)?', s)
    if not m:
        return None
    try:
        val = float(m.group(1))
    except (ValueError, TypeError):
        return None
    unit = (m.group(2) or '').lower()
    if unit.startswith('h'):
        return val * 60
    return val

def _sum_table_duration(headers, rows):
    """Fuente UNICA de verdad para el calculo de duracion de una tabla.
    Detecta la columna 'Duracion'/'Tiempo' y suma sus minutos, EXCLUYENDO filas
    total/subtotal y filas 'solo-duracion' (todas las demas celdas vacias).
    La usan tanto el Word (_add_styled_table) como el TOTAL del markdown
    (ai_helper._normalizar_total_carta), garantizando que ambos coincidan.
    Devuelve (duration_col, total_min, total_row_indices, any_numeric)."""
    duration_col = None
    for idx, h in enumerate(headers):
        if any(k in str(h).lower() for k in ['duraci', 'tiempo']):
            duration_col = idx
            break
    if duration_col is None:
        return None, 0.0, [], False
    total = 0.0
    any_numeric = False
    total_row_indices = []
    for ridx, r in enumerate(rows):
        row_label = " ".join(str(c) for c in r[:max(1, duration_col)]).strip().lower()
        non_dur = [str(c).strip() for idx, c in enumerate(r) if idx != duration_col]
        is_summary_only = (duration_col < len(r) and str(r[duration_col]).strip() != "" and not any(non_dur))
        is_total_row = bool(re.match(r'^\s*(sub)?total\b', row_label)) or row_label.startswith('total ') or is_summary_only
        if is_total_row:
            total_row_indices.append(ridx)
            continue
        if duration_col < len(r):
            v = _extract_minutes(r[duration_col])
            if v is not None:
                total += v
                any_numeric = True
    return duration_col, total, total_row_indices, any_numeric

def _add_styled_table(doc, headers, rows, sc, duration_tracker=None):
    duration_col, total, total_row_indices, any_numeric = _sum_table_duration(headers, rows)

    subtotal_row = None
    subtotal_minutes = 0.0
    if duration_col is not None and rows:
        if any_numeric and total > 0:
            total_int = int(total) if total == int(total) else round(total, 1)
            for ridx in total_row_indices:
                if duration_col < len(rows[ridx]):
                    rows[ridx] = list(rows[ridx])
                    rows[ridx][duration_col] = f"{total_int} min"
            subtotal_minutes = total
            subtotal_row = [""] * len(headers)
            label_col = max(0, duration_col - 1)
            subtotal_row[label_col] = "SUBTOTAL"
            subtotal_row[duration_col] = f"{total_int} min"
            if duration_tracker is not None:
                duration_tracker.append(subtotal_minutes)

    total_rows = 1 + len(rows) + (1 if subtotal_row else 0)
    table = doc.add_table(rows=total_rows, cols=len(headers))
    table.alignment = WD_TABLE_ALIGNMENT.CENTER
    table.style = 'Table Grid'
    table.autofit = False

    col_widths = _compute_column_widths(headers)
    for col_idx, width in enumerate(col_widths):
        for row in table.rows:
            if col_idx < len(row.cells):
                row.cells[col_idx].width = width

    # Con tblLayout="fixed", Word usa el tblGrid (no el ancho de celda
    # individual) para dibujar las columnas. Fijar cell.width por sí solo
    # no actualiza el tblGrid — hay que hacerlo explícitamente o el ancho
    # visual no cambia aunque el XML de las celdas sí lo tenga.
    _tbl = table._tbl
    _tblGrid = _tbl.find(qn('w:tblGrid'))
    if _tblGrid is not None:
        _grid_cols = _tblGrid.findall(qn('w:gridCol'))
        for col_idx, width in enumerate(col_widths):
            if col_idx < len(_grid_cols):
                _grid_cols[col_idx].set(qn('w:w'), str(width.twips))

    header_row = table.rows[0]
    for i, header_text in enumerate(headers):
        cell = header_row.cells[i]
        _set_cell_text(cell, header_text, bold=True, size=9, color=sc["table_header_color"], align=WD_ALIGN_PARAGRAPH.CENTER, font_name=sc["font_body"])
        _set_cell_shading(cell, sc["table_header_bg"])

    for row_idx, row_data in enumerate(rows):
        row = table.rows[row_idx + 1]
        bg = sc["table_row_alt"] if row_idx % 2 == 0 else "FFFFFF"
        for col_idx, cell_text in enumerate(row_data):
            cell = row.cells[col_idx]
            _set_cell_text(cell, str(cell_text), size=9, color=sc["body_color"], font_name=sc["font_body"])
            _set_cell_shading(cell, bg)

    if subtotal_row:
        srow = table.rows[-1]
        for col_idx, cell_text in enumerate(subtotal_row):
            cell = srow.cells[col_idx]
            _set_cell_text(cell, str(cell_text), bold=True, size=9, color=sc["heading_color"], font_name=sc["font_body"], align=WD_ALIGN_PARAGRAPH.CENTER)
            _set_cell_shading(cell, "E8E8E8")

    return table

import logging
import tempfile
import urllib.request
import urllib.parse
import json as _json

_pptx_logger = logging.getLogger('pptx_generator')

_pexels_usage_log = []
_last_pexels_credits = []

_BROWSER_UA = ("Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 "
               "(KHTML, like Gecko) Chrome/124.0.0.0 Safari/537.36")

def pop_last_pexels_usage():
    global _last_pexels_credits
    items = _last_pexels_credits
    _last_pexels_credits = []
    return items

def _apply_fill_alpha(shape, alpha_val):
    try:
        from pptx.oxml.ns import qn as _qn
        from lxml import etree as _etree
        sf = shape._element.spPr.find(_qn('a:solidFill'))
        if sf is None:
            return
        srgb = sf.find(_qn('a:srgbClr'))
        if srgb is None:
            return
        a = _etree.SubElement(srgb, _qn('a:alpha'))
        a.set('val', str(alpha_val))
    except Exception:
        pass

def _fetch_pexels_image(query, orientation="landscape", size="medium", session_ctx=None):
    api_key = os.environ.get("PEXELS_API_KEY", "")
    if not api_key:
        return None
    if session_ctx and session_ctx.get("rate_paused"):
        return None
    try:
        safe_query = urllib.parse.quote(query)
        url = f"https://api.pexels.com/v1/search?query={safe_query}&per_page=3&orientation={orientation}&size={size}"
        req = urllib.request.Request(url, headers={
            "Authorization": api_key,
            "User-Agent": _BROWSER_UA,
            "Accept": "application/json",
        })
        with urllib.request.urlopen(req, timeout=8) as resp:
            remaining = resp.headers.get("X-Ratelimit-Remaining", "")
            if remaining.isdigit() and int(remaining) < 10:
                if session_ctx is not None:
                    session_ctx["rate_paused"] = True
                _pptx_logger.warning(f"Pexels rate-limit low ({remaining} remaining) — pausing image requests")
            data = _json.loads(resp.read().decode())
        photos = data.get("photos", [])
        if not photos:
            return None
        photo = photos[0]
        img_url = photo.get("src", {}).get("large", photo.get("src", {}).get("medium", ""))
        if not img_url:
            return None
        tmp = tempfile.NamedTemporaryFile(suffix=".jpg", delete=False, dir=tempfile.gettempdir())
        img_req = urllib.request.Request(img_url, headers={"User-Agent": _BROWSER_UA})
        with urllib.request.urlopen(img_req, timeout=12) as img_resp, open(tmp.name, "wb") as out_f:
            out_f.write(img_resp.read())
        entry = {
            "query": query,
            "photo_id": photo.get("id"),
            "photographer": photo.get("photographer", ""),
            "photographer_url": photo.get("photographer_url", ""),
            "url": photo.get("url", ""),
        }
        _pexels_usage_log.append(entry)
        if len(_pexels_usage_log) > 500:
            _pexels_usage_log.pop(0)
        if session_ctx is not None:
            session_ctx["credits"].append(entry)
        return tmp.name
    except urllib.error.HTTPError as e:
        if e.code == 429:
            if session_ctx is not None:
                session_ctx["rate_paused"] = True
            _pptx_logger.warning("Pexels 429 Too Many Requests — pausing for this session")
        else:
            _pptx_logger.warning(f"Pexels HTTP {e.code} for '{query}': {e}")
        return None
    except Exception as e:
        _pptx_logger.warning(f"Pexels image fetch failed for '{query}': {e}")
        return None

def _extract_image_keywords(course_name, slide_title=""):
    stop_words = {'de', 'del', 'la', 'el', 'los', 'las', 'en', 'y', 'a', 'para', 'por', 'con', 'que',
                  'un', 'una', 'su', 'al', 'como', 'se', 'es', 'no', 'más', 'pero', 'o', 'si',
                  'técnica', 'expositiva', 'demostrativa', 'diálogo', 'discusión', 'desarrollo',
                  'apertura', 'cierre', 'encuadre', 'slide', 'diapositiva', 'cont', 'módulo',
                  'evaluación', 'sumativa', 'formativa', 'diagnóstica', 'instrucciones'}
    combined = f"{course_name} {slide_title}"
    words = [w.strip('()[]¿?¡!:.,;') for w in combined.lower().split() if len(w) > 2]
    keywords = [w for w in words if w not in stop_words][:5]
    return " ".join(keywords) if keywords else course_name

def get_pexels_usage_stats():
    return {
        "total_fetches": len(_pexels_usage_log),
        "recent": _pexels_usage_log[-10:] if _pexels_usage_log else [],
        "has_api_key": bool(os.environ.get("PEXELS_API_KEY", "")),
    }

_PPTX_PALETTES = {
    "industrial": {
        "bg_dark": (30, 40, 55), "bg_section": (20, 30, 48),
        "accent": (0, 150, 199), "accent2": (0, 188, 212),
        "title": (255, 255, 255), "body": (230, 235, 240),
        "subtitle": (160, 180, 200), "muted": (120, 140, 160),
    },
    "salud": {
        "bg_dark": (25, 45, 35), "bg_section": (18, 38, 28),
        "accent": (46, 204, 113), "accent2": (39, 174, 96),
        "title": (255, 255, 255), "body": (230, 240, 235),
        "subtitle": (160, 200, 180), "muted": (120, 160, 140),
    },
    "educacion": {
        "bg_dark": (35, 30, 55), "bg_section": (28, 22, 48),
        "accent": (124, 58, 237), "accent2": (139, 92, 246),
        "title": (255, 255, 255), "body": (235, 230, 245),
        "subtitle": (180, 170, 210), "muted": (140, 130, 170),
    },
    "corporativo": {
        "bg_dark": (32, 32, 38), "bg_section": (25, 25, 32),
        "accent": (250, 128, 114), "accent2": (255, 160, 148),
        "title": (255, 255, 255), "body": (235, 235, 240),
        "subtitle": (180, 180, 195), "muted": (140, 140, 155),
    },
    "seguridad": {
        "bg_dark": (45, 30, 15), "bg_section": (38, 24, 10),
        "accent": (255, 165, 0), "accent2": (255, 193, 7),
        "title": (255, 255, 255), "body": (240, 235, 225),
        "subtitle": (200, 185, 160), "muted": (165, 150, 130),
    },
}

def _detect_course_palette(course_name):
    cn = course_name.lower()
    if any(k in cn for k in ['lubric', 'soldad', 'manufactur', 'maquina', 'industrial', 'hidraul',
                               'mecan', 'electric', 'neumat', 'metalurg', 'mantenimiento', 'producci']):
        return "industrial"
    if any(k in cn for k in ['salud', 'medic', 'enferm', 'hospital', 'clinic', 'farmac',
                               'nutrici', 'higiene', 'sanitar', 'primer auxilio']):
        return "salud"
    if any(k in cn for k in ['pedagog', 'didact', 'enseñ', 'educaci', 'docen', 'capacitaci',
                               'formaci', 'aprend', 'tutor']):
        return "educacion"
    if any(k in cn for k in ['seguridad', 'proteccion', 'riesgo', 'emergencia', 'incendio',
                               'evacuaci', 'epp', 'norma oficial']):
        return "seguridad"
    return "corporativo"

_SLIDE_ICON_MAP = [
    (['comprobacion', 'verificacion', 'recursos y condiciones', 'momento cero', 'checklist', 'lista de verificacion'], '\u2705'),
    (['rompe hielo', 'rompehielo', 'tarjeta', 'integracion', 'presentacion de los participantes'], '\U0001F91D'),
    (['practica', 'ejercicio', 'actividad practica', 'caso practico', 'taller'], '\U0001F4DC'),
    (['descanso', 'receso', 'pausa'], '\u2615'),
    (['energizante', 'semaforo', 'energia', 'dinamica grupal', 'actividad energizante'], '\u26A1'),
    (['reflexion', 'reflexiona', 'expectativa', 'compromiso', 'aplicacion'], '\U0001F4AD'),
    (['temario', 'ruta del curso', 'programa', 'indice', 'agenda', 'contenido tematico'], '\U0001F4D1'),
    (['objetivo', 'proposito', 'meta del curso'], '\U0001F4CB'),
    (['beneficio'], '\u2705'),
    (['evaluacion', 'diagnostica', 'formativa', 'sumativa', 'ponderacion', 'satisfaccion', 'encuesta', 'retroalimentacion'], '\U0001F4CA'),
    (['acuerdo', 'convivencia', 'regla', 'norma de convivencia'], '\U0001F4DC'),
    (['resumen', 'conclusion', 'sintesis', 'cierre', 'gracias', 'agradecimiento', 'despedida'], '\U0001F3C1'),
    (['bienvenid', 'encuadre', 'apertura', 'inicio de la sesion'], '\U0001F3AF'),
    (['bibliografia', 'referencia', 'fuente', 'continuidad', 'sugerencia'], '\U0001F4DA'),
    (['desarrollo', 'tema', 'tecnica', 'modulo', 'contenido'], '\u2699\uFE0F'),
]

def _get_slide_icon(title_text):
    import unicodedata
    t = ''.join(c for c in unicodedata.normalize('NFKD', title_text.lower()) if not unicodedata.category(c).startswith('M'))
    for keywords, icon in _SLIDE_ICON_MAP:
        if any(kw in t for kw in keywords):
            return icon
    return ''

def _should_use_table_layout(title, content_lines):
    import unicodedata
    t = ''.join(c for c in unicodedata.normalize('NFKD', title.lower()) if not unicodedata.category(c).startswith('M'))
    has_keyword = any(k in t for k in ['objetivo', 'evaluacion', 'ponderacion', 'verificacion', 'criterio'])
    has_table = sum(1 for line in content_lines if line.count('|') >= 2) >= 2
    return has_keyword and has_table

def _should_use_flow_layout(title, content_lines):
    import unicodedata
    t = ''.join(c for c in unicodedata.normalize('NFKD', title.lower()) if not unicodedata.category(c).startswith('M'))
    has_keyword = any(k in t for k in ['proceso', 'flujo', 'secuencia', 'fases', 'etapas', 'estructura del curso'])
    bullets = [l.strip() for l in content_lines if l.strip() and (l.strip().startswith('-') or l.strip().startswith('*') or re.match(r'^\d+[\.\)]\s', l.strip()))]
    short_enough = all(len(re.sub(r'^[\-\*\d\.\)]+\s*', '', b)) < 45 for b in bullets) if bullets else False
    return has_keyword and 3 <= len(bullets) <= 6 and short_enough

def generate_pptx_from_slides(slide_data, course_name="", modalidad="facilitacion", logo_path=None, curso_datos=None):
    from pptx import Presentation as PptxPresentation
    from pptx.util import Inches as PptxInches, Pt as PptxPt, Emu
    from pptx.dml.color import RGBColor as PptxRGB
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.enum.shapes import MSO_SHAPE

    global _last_pexels_credits
    _last_pexels_credits = []

    os.makedirs(GENERATED_DIR, exist_ok=True)
    prs = PptxPresentation()
    prs.slide_width = PptxInches(13.333)
    prs.slide_height = PptxInches(7.5)

    curso_datos = curso_datos or {}

    _es_ejecutiva = (modalidad == "ejecutiva")
    MAX_BULLETS_PER_SLIDE = 6 if _es_ejecutiva else 8

    palette_name = _detect_course_palette(course_name)
    pal = _PPTX_PALETTES[palette_name]

    def _clamp(v):
        return max(0, min(255, int(v)))
    def _lighten(rgb, dr, dg, db):
        return (_clamp(rgb[0] + dr), _clamp(rgb[1] + dg), _clamp(rgb[2] + db))
    def _darken(rgb, f):
        return (_clamp(rgb[0] * f), _clamp(rgb[1] * f), _clamp(rgb[2] * f))

    _bg = pal["bg_dark"]
    BG_DARK = PptxRGB(*_bg)
    BG_OUTRO = PptxRGB(*_darken(_bg, 0.78))
    BG_CARD = PptxRGB(*_lighten(_bg, 12, 12, 16))
    BG_ALT = PptxRGB(*_lighten(_bg, 24, 24, 34))
    BG_SECTION = PptxRGB(*pal["bg_section"])
    ACCENT = PptxRGB(*pal["accent"])
    ACCENT2 = PptxRGB(*pal["accent2"])
    ACCENT_DARK = PptxRGB(*_darken(pal["accent"], 0.8))
    TITLE_CLR = PptxRGB(*pal["title"])
    BODY_CLR = PptxRGB(*pal["body"])
    SUBTITLE_CLR = PptxRGB(*pal["subtitle"])
    MUTED_CLR = PptxRGB(*pal["muted"])
    FONT = "Calibri"

    page_counter = [0]
    course_short = _sanitize_text(course_name)[:46] if course_name else "Curso de Capacitación"

    def _set_bg(slide, color):
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = color

    def _blank(bg=None):
        s = prs.slides.add_slide(prs.slide_layouts[6])
        _set_bg(s, bg if bg is not None else BG_DARK)
        return s

    def _rect(slide, l, t, w, h, color):
        sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, PptxInches(l), PptxInches(t), PptxInches(w), PptxInches(h))
        sh.fill.solid()
        sh.fill.fore_color.rgb = color
        sh.line.fill.background()
        try:
            sh.shadow.inherit = False
        except Exception:
            pass
        return sh

    def _text(slide, l, t, w, h, content, size=14, color=None, bold=False, align=PP_ALIGN.LEFT, anchor=None):
        box = slide.shapes.add_textbox(PptxInches(l), PptxInches(t), PptxInches(w), PptxInches(h))
        tf = box.text_frame
        tf.word_wrap = True
        try:
            tf.margin_left = PptxInches(0.04)
            tf.margin_right = PptxInches(0.04)
            tf.margin_top = PptxInches(0.02)
            tf.margin_bottom = PptxInches(0.02)
        except Exception:
            pass
        if anchor is not None:
            try:
                tf.vertical_anchor = anchor
            except Exception:
                pass
        p = tf.paragraphs[0]
        p.alignment = align
        runs = content if isinstance(content, list) else [(content, size, color if color is not None else BODY_CLR, bold)]
        for (txt, rs, rc, rb) in runs:
            r = p.add_run()
            r.text = txt
            r.font.size = PptxPt(rs)
            r.font.bold = rb
            r.font.color.rgb = rc
            r.font.name = FONT
        return box, tf

    def _norm(s):
        import unicodedata
        return ''.join(c for c in unicodedata.normalize('NFKD', (s or "").lower()) if not unicodedata.category(c).startswith('M'))

    def _clean_bullets(lines):
        out = []
        for line in (lines or []):
            s = _sanitize_text(line.strip())
            if not s:
                continue
            is_num = bool(re.match(r'^\d+[\.\)]\s', s))
            is_dash = s.startswith('-') or s.startswith('*')
            is_bullet = is_num or is_dash
            if is_dash:
                clean = re.sub(r'^[\-\*]+\s*', '', s)
            elif is_num:
                clean = re.sub(r'^\d+[\.\)]\s*', '', s)
            else:
                clean = s
            is_sub = (not is_bullet) and clean.endswith(':') and len(clean) < 60
            out.append((clean, is_sub))
        return out

    def _extract_pairs(lines):
        pairs = []
        for line in (lines or []):
            s = _sanitize_text(line.strip())
            if not s:
                continue
            s = re.sub(r'^[\-\*]+\s*', '', s)
            if ':' in s:
                k, v = s.split(':', 1)
                k = k.strip(); v = v.strip()
                if k and v and len(k) <= 34:
                    pairs.append((k, v))
                else:
                    return []
            else:
                return []
        return pairs

    def _notes(slide, notes_text):
        if notes_text:
            try:
                slide.notes_slide.notes_text_frame.text = _sanitize_text(notes_text)
            except Exception:
                pass

    def _header(slide, title_text, numbered=True):
        _rect(slide, 0.0, 0.0, 13.333, 1.55, BG_CARD)
        _rect(slide, 0.4, 1.5, 12.53, 0.04, ACCENT)
        _rect(slide, 0.4, 0.2, 0.1, 1.15, ACCENT)
        icon = _get_slide_icon(title_text)
        disp = f"{icon}  {_sanitize_text(title_text)}" if icon else _sanitize_text(title_text)
        _text(slide, 0.6, 0.18, 11.6, 1.2, disp, size=22, color=TITLE_CLR, bold=True, anchor=MSO_ANCHOR.MIDDLE)
        if numbered:
            page_counter[0] += 1
            _text(slide, 12.45, 0.05, 0.78, 0.35, str(page_counter[0]), size=9, color=MUTED_CLR, align=PP_ALIGN.RIGHT)

    def _footer(slide):
        _rect(slide, 0.4, 7.05, 12.53, 0.02, BG_ALT)
        _logo_ok = False
        if logo_path:
            try:
                if os.path.exists(logo_path):
                    slide.shapes.add_picture(logo_path, PptxInches(0.4), PptxInches(7.0), height=PptxInches(0.4))
                    _logo_ok = True
            except Exception:
                _logo_ok = False
        if not _logo_ok:
            _text(slide, 0.4, 7.08, 3.0, 0.35, "[ Su logo ]", size=8, color=MUTED_CLR, align=PP_ALIGN.LEFT)
        _text(slide, 3.5, 7.08, 6.33, 0.35, course_short, size=8, color=MUTED_CLR, align=PP_ALIGN.CENTER)

    _add_footer = _footer

    def _add_accent_line(slide, left, top, width, height=0.05):
        line = slide.shapes.add_shape(1, PptxInches(left), PptxInches(top), PptxInches(width), PptxInches(height))
        line.fill.solid()
        line.fill.fore_color.rgb = ACCENT
        line.line.fill.background()

    _pexels_session = {"rate_paused": False, "credits": []}
    images_inserted = [0]
    MAX_IMAGES = 4
    _temp_image_files = []

    def _try_add_image(slide, query, left, top, width, height):
        if images_inserted[0] >= MAX_IMAGES:
            return False
        img_path = _fetch_pexels_image(query, session_ctx=_pexels_session)
        if not img_path:
            return False
        try:
            slide.shapes.add_picture(img_path, PptxInches(left), PptxInches(top), PptxInches(width), PptxInches(height))
            images_inserted[0] += 1
            _temp_image_files.append(img_path)
            return True
        except Exception:
            try:
                os.unlink(img_path)
            except Exception:
                pass
            return False

    def _cleanup_temp_images():
        for f in _temp_image_files:
            try:
                os.unlink(f)
            except Exception:
                pass

    def _add_title_slide():
        slide = _blank(BG_DARK)
        _rect(slide, 10.33, 0.0, 3.0, 2.2, ACCENT_DARK)
        _rect(slide, 11.33, 0.0, 2.0, 1.6, ACCENT)
        _rect(slide, 0.0, 5.5, 2.5, 2.0, ACCENT_DARK)
        _rect(slide, 0.0, 6.2, 1.8, 1.3, ACCENT)
        _rect(slide, 0.8, 2.2, 8.5, 0.04, ACCENT)
        _rect(slide, 0.8, 5.3, 8.5, 0.04, ACCENT)

        _text(slide, 0.8, 1.55, 10.0, 0.55, "CURSO-TALLER DE CAPACITACIÓN PRESENCIAL GRUPAL",
              size=11, color=ACCENT, bold=True)
        _text(slide, 0.8, 2.35, 9.2, 1.7, _sanitize_text(course_name if course_name else "Presentación del Curso"),
              size=24, color=TITLE_CLR, bold=True)

        instructor = curso_datos.get('instructor') or "[ Nombre del Instructor ]"
        _text(slide, 0.8, 4.15, 9.0, 0.5, _sanitize_text("IMPARTIDO POR " + instructor),
              size=16, color=ACCENT, bold=True)

        _rect(slide, 0.8, 4.75, 11.73, 1.6, BG_CARD)
        periodo = curso_datos.get('periodo') or curso_datos.get('fecha') or "[ Fecha ]"
        horario = curso_datos.get('horario')
        if horario:
            periodo = f"{periodo} · {horario}"
        sede = curso_datos.get('sede') or "[ Sede ]"
        duracion = curso_datos.get('duracion') or "[ Duración ]"
        participantes = curso_datos.get('participantes') or "[ No. de participantes ]"
        _text(slide, 0.98, 4.9, 5.6, 0.48, _sanitize_text(f"\U0001F4C5  Periodo: {periodo}"), size=11, color=BODY_CLR, anchor=MSO_ANCHOR.MIDDLE)
        _text(slide, 0.98, 5.45, 5.6, 0.48, _sanitize_text(f"\U0001F4CD  Sede: {sede}"), size=11, color=BODY_CLR, anchor=MSO_ANCHOR.MIDDLE)
        _text(slide, 6.7, 4.9, 5.6, 0.48, _sanitize_text(f"\u23F0  Duración: {duracion}"), size=11, color=BODY_CLR, anchor=MSO_ANCHOR.MIDDLE)
        _text(slide, 6.7, 5.45, 5.6, 0.48, _sanitize_text(f"\U0001F465  Participantes: {participantes}"), size=11, color=BODY_CLR, anchor=MSO_ANCHOR.MIDDLE)

        _cover_logo_ok = False
        if logo_path:
            try:
                if os.path.exists(logo_path):
                    slide.shapes.add_picture(logo_path, PptxInches(9.8), PptxInches(1.0), height=PptxInches(0.9))
                    _cover_logo_ok = True
            except Exception:
                _cover_logo_ok = False
        if not _cover_logo_ok:
            _rect(slide, 9.8, 1.0, 2.8, 0.9, BG_ALT)
            _text(slide, 9.8, 1.0, 2.8, 0.9, "[ LOGO ]", size=11, color=MUTED_CLR, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

        disenador = curso_datos.get('disenador')
        footer_txt = "BIENVENIDOS AL CURSO"
        if disenador:
            footer_txt += f"  ·  Diseño: {disenador}"
        _text(slide, 0.8, 7.12, 12.0, 0.32, _sanitize_text(footer_txt), size=8, color=MUTED_CLR)

    def _add_section_slide(title_text, notes_text=""):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        _set_bg(slide, BG_SECTION)

        sec_query = _extract_image_keywords(course_name, title_text)
        img_added = _try_add_image(slide, sec_query, 8.5, 1.5, 4.3, 4.5)

        _add_accent_line(slide, 1.5, 3.0, 3.5, 0.08)

        icon = _get_slide_icon(title_text)
        display_title = f"{icon}  {_sanitize_text(title_text)}" if icon else _sanitize_text(title_text)

        title_width = 6.5 if img_added else 10.3
        title_box = slide.shapes.add_textbox(PptxInches(1.5), PptxInches(3.3), PptxInches(title_width), PptxInches(1.5))
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = display_title
        p.font.size = PptxPt(36)
        p.font.bold = True
        p.font.color.rgb = ACCENT
        p.alignment = PP_ALIGN.LEFT

        if notes_text:
            slide.notes_slide.notes_text_frame.text = _sanitize_text(notes_text)
        _add_footer(slide)

    def _render_bullets_card(slide, body_lines, top=1.68, height=5.2):
        _rect(slide, 0.4, top, 12.53, height, BG_CARD)
        items = _clean_bullets(body_lines)
        if not items:
            return
        box = slide.shapes.add_textbox(PptxInches(0.6), PptxInches(top + 0.17), PptxInches(12.13), PptxInches(height - 0.32))
        tf = box.text_frame
        tf.word_wrap = True
        first = True
        for text, is_sub in items:
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            first = False
            p.space_after = PptxPt(8)
            if is_sub:
                r = p.add_run()
                r.text = text
                r.font.size = PptxPt(16)
                r.font.bold = True
                r.font.color.rgb = ACCENT
                r.font.name = FONT
                p.space_before = PptxPt(10)
            else:
                m = p.add_run()
                m.text = "\u25B8  "
                m.font.size = PptxPt(14)
                m.font.bold = True
                m.font.color.rgb = ACCENT
                m.font.name = FONT
                r = p.add_run()
                r.text = text
                r.font.size = PptxPt(15)
                r.font.color.rgb = BODY_CLR
                r.font.name = FONT

    def _add_content_slide(title_text, body_lines, notes_text=""):
        slide = _blank(BG_DARK)
        _header(slide, title_text)
        _render_bullets_card(slide, body_lines)
        _footer(slide)
        _notes(slide, notes_text)

    def _add_two_column_slide(title_text, body_lines, notes_text=""):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        _set_bg(slide, BG_DARK)

        title_box = slide.shapes.add_textbox(PptxInches(0.8), PptxInches(0.4), PptxInches(11.7), PptxInches(1.0))
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = _sanitize_text(title_text)
        p.font.size = PptxPt(28)
        p.font.bold = True
        p.font.color.rgb = ACCENT
        p.alignment = PP_ALIGN.LEFT

        _add_accent_line(slide, 0.8, 1.4, 2.0)

        left_lines = []
        right_lines = []
        current_col = left_lines
        current_header = None
        for line in body_lines:
            stripped = line.strip()
            is_bullet = stripped.startswith("-") or stripped.startswith("*") or re.match(r'^\d+[\.\)]\s', stripped)
            if not is_bullet and stripped and not stripped.startswith("#"):
                if current_header is not None and current_col is left_lines:
                    current_col = right_lines
                current_header = stripped
                current_col.append(stripped)
            else:
                current_col.append(stripped)

        if not right_lines:
            mid = len(body_lines) // 2
            left_lines = body_lines[:mid]
            right_lines = body_lines[mid:]

        for col_lines, left_pos in [(left_lines, 0.8), (right_lines, 6.8)]:
            col_box = slide.shapes.add_textbox(PptxInches(left_pos), PptxInches(1.7), PptxInches(5.5), PptxInches(4.8))
            ctf = col_box.text_frame
            ctf.word_wrap = True
            for i, line in enumerate(col_lines):
                line = _sanitize_text(line.strip())
                if not line:
                    continue
                cp = ctf.paragraphs[0] if i == 0 else ctf.add_paragraph()
                is_bullet = line.startswith("-") or line.startswith("*") or re.match(r'^\d+[\.\)]\s', line)
                is_subheading = not is_bullet and line.endswith(":") and len(line) < 60
                clean_line = re.sub(r'^[\-\*\d\.\)]+\s*', '', line) if is_bullet else line
                cp.text = clean_line
                if is_subheading:
                    cp.font.size = PptxPt(18)
                    cp.font.bold = True
                    cp.font.color.rgb = ACCENT2
                    cp.space_before = PptxPt(8)
                elif is_bullet:
                    cp.font.size = PptxPt(15)
                    cp.font.color.rgb = BODY_CLR
                    cp.level = 1
                else:
                    cp.font.size = PptxPt(16)
                    cp.font.color.rgb = BODY_CLR
                cp.space_after = PptxPt(5)

        divider = slide.shapes.add_shape(1, PptxInches(6.5), PptxInches(1.7), PptxInches(0.03), PptxInches(4.8))
        divider.fill.solid()
        divider.fill.fore_color.rgb = MUTED_CLR
        divider.line.fill.background()

        if notes_text:
            slide.notes_slide.notes_text_frame.text = _sanitize_text(notes_text)
        _add_footer(slide)

    def _build_pexels_attribution():
        if not _pexels_session["credits"]:
            return []
        seen = set()
        lines = []
        for entry in _pexels_session["credits"]:
            photographer = entry.get("photographer", "")
            if photographer and photographer not in seen:
                seen.add(photographer)
                lines.append(f"Foto: {photographer} — Pexels (pexels.com)")
        return lines

    def _add_closing_slide(title_text="¡Gracias por su participación!", body_lines=None, notes_text=""):
        slide = _blank(BG_DARK)
        _header(slide, title_text or "¡Gracias por su participación!")
        items = body_lines if (body_lines and any(_sanitize_text(b.strip()) for b in body_lines)) else [
            "Gracias por su participación activa durante todo el curso",
            "Apliquen lo aprendido en su entorno profesional",
            "El verdadero aprendizaje comienza ahora",
        ]
        _render_bullets_card(slide, items)
        _footer(slide)
        _notes(slide, notes_text)

    def _add_numbered_slide(title_text, items, notes_text=""):
        slide = _blank(BG_DARK)
        _header(slide, title_text)
        n = len(items)
        top = 1.68
        bottom = 6.9
        gap = 0.06
        row_h = min(0.92, (bottom - top - (n - 1) * gap) / n) if n > 0 else 0.92
        if row_h < 0.4:
            row_h = 0.4
        step = row_h + gap
        for i, txt in enumerate(items):
            y = top + i * step
            _rect(slide, 0.4, y, 12.53, row_h, BG_ALT)
            _text(slide, 0.52, y, 0.6, row_h, str(i + 1), size=20, color=ACCENT, bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
            _rect(slide, 1.18, y + row_h * 0.12, 0.03, row_h * 0.76, ACCENT_DARK)
            _text(slide, 1.3, y, 11.4, row_h, _sanitize_text(txt), size=13, color=TITLE_CLR, anchor=MSO_ANCHOR.MIDDLE)
        _footer(slide)
        _notes(slide, notes_text)

    def _add_keyvalue_slide(title_text, pairs, notes_text=""):
        slide = _blank(BG_DARK)
        _header(slide, title_text)
        n = len(pairs)
        top = 1.7
        bottom = 6.9
        gap = 0.07
        row_h = min(0.72, (bottom - top - (n - 1) * gap) / n) if n > 0 else 0.72
        if row_h < 0.4:
            row_h = 0.4
        step = row_h + gap
        for i, (label, value) in enumerate(pairs):
            y = top + i * step
            _rect(slide, 0.4, y, 12.53, row_h, BG_ALT if i % 2 == 0 else BG_CARD)
            _text(slide, 0.55, y, 3.25, row_h, _sanitize_text(label), size=11, color=ACCENT, bold=True, anchor=MSO_ANCHOR.MIDDLE)
            _text(slide, 3.9, y, 8.6, row_h, _sanitize_text(value), size=12, color=TITLE_CLR, anchor=MSO_ANCHOR.MIDDLE)
        _footer(slide)
        _notes(slide, notes_text)

    def _add_break_slide(title_text, body_lines, notes_text=""):
        slide = _blank(BG_DARK)
        _footer(slide)
        _rect(slide, 2.5, 1.8, 8.33, 3.9, BG_CARD)
        _rect(slide, 2.5, 1.8, 8.33, 0.08, ACCENT)
        _rect(slide, 2.5, 5.62, 8.33, 0.08, ACCENT)
        _text(slide, 2.5, 2.0, 8.33, 1.3, "\u2615", size=64, color=ACCENT, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        _text(slide, 2.5, 3.25, 8.33, 0.8, "DESCANSO", size=36, color=TITLE_CLR, bold=True, align=PP_ALIGN.CENTER)
        sub = ""
        for l in (body_lines or []):
            ls = _sanitize_text(re.sub(r'^[\-\*\d\.\)]+\s*', '', l.strip()))
            if ls:
                sub = ls
                break
        if not sub:
            sub = "Regreso puntual para continuar"
        _text(slide, 2.5, 4.15, 8.33, 0.6, sub, size=14, color=BODY_CLR, align=PP_ALIGN.CENTER)
        page_counter[0] += 1
        _text(slide, 2.5, 4.85, 8.33, 0.4, str(page_counter[0]), size=11, color=MUTED_CLR, align=PP_ALIGN.CENTER)
        _notes(slide, notes_text)

    def _add_semaforo_slide(title_text, body_lines, notes_text=""):
        slide = _blank(BG_DARK)
        _header(slide, title_text)
        colors = [("\U0001F7E2", "VERDE", (46, 125, 50)), ("\U0001F7E1", "AMARILLO", (245, 127, 23)), ("\U0001F534", "ROJO", (198, 40, 40))]
        descs = {"verde": "", "amarillo": "", "rojo": ""}
        for l in (body_lines or []):
            ls = _norm(l)
            clean = _sanitize_text(re.sub(r'^[\-\*\d\.\)]+\s*', '', l.strip()))
            for key in descs:
                if key in ls and not descs[key]:
                    d = re.sub(r'(?i)^[^:]*' + key + r'[^:]*:\s*', '', clean)
                    if d == clean:
                        d = re.sub(r'(?i)\b' + key + r'\b[\s:\u2013\-]*', '', clean).strip()
                    descs[key] = d
        defaults = {"verde": "La decisión procede tal como está", "amarillo": "Procede, pero requiere ajustes", "rojo": "No procede — requiere revisión"}
        top = 1.7
        row_h = 1.54
        gap = 0.12
        for i, (emoji, name, rgb) in enumerate(colors):
            y = top + i * (row_h + gap)
            _rect(slide, 0.4, y, 12.53, row_h, BG_CARD)
            _rect(slide, 0.4, y, 0.18, row_h, PptxRGB(*rgb))
            _text(slide, 0.65, y, 1.2, row_h, emoji, size=30, color=TITLE_CLR, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
            _text(slide, 2.0, y + 0.12, 4.0, 0.55, name, size=18, color=PptxRGB(*rgb), bold=True)
            key = name.lower()
            _text(slide, 2.0, y + 0.7, 10.35, 0.7, descs.get(key) or defaults[key], size=13, color=BODY_CLR)
        _footer(slide)
        _notes(slide, notes_text)

    def _add_brand_outro():
        slide = _blank(BG_OUTRO)
        _rect(slide, 4.67, 2.3, 4.0, 0.05, ACCENT)
        _brand_logo_ok = False
        if logo_path:
            try:
                if os.path.exists(logo_path):
                    _pic = slide.shapes.add_picture(logo_path, PptxInches(5.67), PptxInches(2.7), height=PptxInches(1.2))
                    _pic.left = int((prs.slide_width - _pic.width) / 2)
                    _brand_logo_ok = True
            except Exception:
                _brand_logo_ok = False
        if not _brand_logo_ok:
            _rect(slide, 5.17, 2.7, 3.0, 1.2, BG_CARD)
            _text(slide, 5.17, 2.7, 3.0, 1.2, "[ LOGO ]", size=12, color=MUTED_CLR, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        _text(slide, 1.0, 4.3, 11.33, 0.8, "¡Gracias!", size=32, color=TITLE_CLR, bold=True, align=PP_ALIGN.CENTER)
        _text(slide, 1.0, 5.2, 11.33, 0.5, course_short, size=14, color=ACCENT, align=PP_ALIGN.CENTER)
        _rect(slide, 4.67, 5.9, 4.0, 0.05, ACCENT)

    def _split_overflow(slide_dict):
        body = slide_dict.get("content", [])
        if len(body) <= MAX_BULLETS_PER_SLIDE:
            return [slide_dict]
        if _es_ejecutiva:
            trimmed = dict(slide_dict)
            trimmed["content"] = body[:MAX_BULLETS_PER_SLIDE]
            return [trimmed]
        chunks = []
        for i in range(0, len(body), MAX_BULLETS_PER_SLIDE):
            chunk = body[i:i + MAX_BULLETS_PER_SLIDE]
            suffix = f" (cont.)" if i > 0 else ""
            chunks.append({
                "title": slide_dict["title"] + suffix,
                "content": chunk,
                "notes": slide_dict.get("notes", "") if i == 0 else "",
                "type": slide_dict.get("type", "content"),
            })
        return chunks

    def _cell_fill(cell, hex_val):
        tcPr = cell._tc.get_or_add_tcPr()
        solidFill = OxmlElement('a:solidFill')
        srgbClr = OxmlElement('a:srgbClr')
        srgbClr.set('val', hex_val)
        solidFill.append(srgbClr)
        tcPr.append(solidFill)

    def _add_table_slide(title_text, body_lines, notes_text=""):
        slide = _blank(BG_DARK)
        _header(slide, title_text)

        accent_hex = '{:02X}{:02X}{:02X}'.format(*_darken(pal["accent"], 0.8))
        card_hex = '{:02X}{:02X}{:02X}'.format(*_lighten(_bg, 12, 12, 16))
        alt_hex = '{:02X}{:02X}{:02X}'.format(*_lighten(_bg, 24, 24, 34))

        table_lines = [l for l in body_lines if l.count('|') >= 2]
        if table_lines:
            headers = [c.strip() for c in table_lines[0].strip('|').split('|') if c.strip()]
            data_lines = []
            for tl in table_lines[1:]:
                if re.match(r'^[\s|:\-]+$', tl.strip()):
                    continue
                cells = [c.strip() for c in tl.strip('|').split('|')]
                if cells:
                    data_lines.append(cells[:len(headers)])

            if headers and data_lines:
                num_rows = min(len(data_lines), 9) + 1
                num_cols = len(headers)
                tbl_shape = slide.shapes.add_table(num_rows, num_cols, PptxInches(0.4), PptxInches(1.8), PptxInches(12.53), PptxInches(4.7))
                tbl = tbl_shape.table

                for ci, h_text in enumerate(headers):
                    cell = tbl.cell(0, ci)
                    cell.text = _sanitize_text(h_text.strip())
                    for paragraph in cell.text_frame.paragraphs:
                        paragraph.alignment = PP_ALIGN.CENTER
                        for run in paragraph.runs:
                            run.font.size = PptxPt(12)
                            run.font.bold = True
                            run.font.color.rgb = TITLE_CLR
                            run.font.name = FONT
                    _cell_fill(cell, accent_hex)

                for ri, row_data in enumerate(data_lines[:9]):
                    for ci, cell_text in enumerate(row_data):
                        if ci >= num_cols:
                            break
                        cell = tbl.cell(ri + 1, ci)
                        cell.text = _sanitize_text(cell_text.strip())
                        for paragraph in cell.text_frame.paragraphs:
                            for run in paragraph.runs:
                                run.font.size = PptxPt(11)
                                run.font.color.rgb = BODY_CLR
                                run.font.name = FONT
                        _cell_fill(cell, alt_hex if ri % 2 == 0 else card_hex)

        _footer(slide)
        _notes(slide, notes_text)

    def _add_flow_slide(title_text, body_lines, notes_text=""):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        _set_bg(slide, BG_DARK)

        vbar = slide.shapes.add_shape(1, PptxInches(0.35), PptxInches(0.4), PptxInches(0.08), PptxInches(1.0))
        vbar.fill.solid()
        vbar.fill.fore_color.rgb = ACCENT
        vbar.line.fill.background()

        icon = _get_slide_icon(title_text)
        display_title = f"{icon}  {_sanitize_text(title_text)}" if icon else _sanitize_text(title_text)

        title_box = slide.shapes.add_textbox(PptxInches(0.8), PptxInches(0.4), PptxInches(11.7), PptxInches(1.0))
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = display_title
        p.font.size = PptxPt(28)
        p.font.bold = True
        p.font.color.rgb = ACCENT
        p.alignment = PP_ALIGN.LEFT

        _add_accent_line(slide, 0.8, 1.4, 2.0)

        items = []
        for line in body_lines:
            stripped = line.strip()
            if stripped and (stripped.startswith('-') or stripped.startswith('*') or re.match(r'^\d+[\.\)]\s', stripped)):
                clean = re.sub(r'^[\-\*\d\.\)]+\s*', '', stripped).strip()
                if clean:
                    items.append(clean)
        if not items:
            items = [l.strip() for l in body_lines if l.strip()]

        n = min(len(items), 6)
        if n == 0:
            _add_content_slide(title_text, body_lines, notes_text)
            return

        total_width = 11.7
        node_width = min(2.2, (total_width - (n - 1) * 0.6) / n)
        arrow_width = 0.4
        total_used = n * node_width + (n - 1) * arrow_width
        start_x = (13.333 - total_used) / 2
        node_y = 3.0
        node_height = 1.4

        for i, item_text in enumerate(items[:n]):
            x = start_x + i * (node_width + arrow_width)
            node = slide.shapes.add_shape(5, PptxInches(x), PptxInches(node_y), PptxInches(node_width), PptxInches(node_height))
            node.fill.solid()
            node.fill.fore_color.rgb = ACCENT if i % 2 == 0 else ACCENT2
            node.line.fill.background()
            ntf = node.text_frame
            ntf.word_wrap = True
            np_ = ntf.paragraphs[0]
            np_.text = item_text[:40]
            np_.font.size = PptxPt(12)
            np_.font.bold = True
            np_.font.color.rgb = TITLE_CLR
            np_.alignment = PP_ALIGN.CENTER

            if i < n - 1:
                arrow_x = x + node_width + 0.02
                arrow_box = slide.shapes.add_textbox(PptxInches(arrow_x), PptxInches(node_y + 0.35), PptxInches(arrow_width), PptxInches(0.5))
                atf = arrow_box.text_frame
                ap = atf.paragraphs[0]
                ap.text = "\u25B6"
                ap.font.size = PptxPt(20)
                ap.font.color.rgb = ACCENT
                ap.alignment = PP_ALIGN.CENTER

        if notes_text:
            slide.notes_slide.notes_text_frame.text = _sanitize_text(notes_text)
        _add_footer(slide)

    _add_title_slide()

    has_closing = False
    for sd in slide_data:
        slide_type = sd.get("type", "content")
        sd_title = sd.get("title", "")
        sd_content = sd.get("content", [])
        sd_notes = sd.get("notes", "")
        nt = _norm(sd_title)

        if slide_type == "closing":
            has_closing = True
            _add_closing_slide(sd_title, sd_content, sd_notes)
            continue

        if any(k in nt for k in ("descanso", "receso", "pausa")):
            _add_break_slide(sd_title, sd_content, sd_notes)
            continue

        if "semaforo" in nt and any("verde" in _norm(l) for l in sd_content) and any("rojo" in _norm(l) for l in sd_content):
            _add_semaforo_slide(sd_title, sd_content, sd_notes)
            continue

        table_lines = [l for l in sd_content if l.count('|') >= 2]
        if len(table_lines) >= 2:
            hdr = [c for c in table_lines[0].strip('|').split('|') if c.strip()]
            if len(hdr) >= 3:
                _add_table_slide(sd_title, sd_content, sd_notes)
                continue
            pairs = []
            for tl in table_lines[1:]:
                if re.match(r'^[\s|:\-]+$', tl.strip()):
                    continue
                cells = [c.strip() for c in tl.strip('|').split('|')]
                if len(cells) >= 2 and (cells[0] or cells[1]):
                    pairs.append((cells[0], cells[1]))
            if pairs:
                _add_keyvalue_slide(sd_title, pairs[:8], sd_notes)
                continue
            _add_content_slide(sd_title, sd_content, sd_notes)
            continue

        non_empty = [l for l in sd_content if l and l.strip()]
        kv_pairs = _extract_pairs(sd_content)
        if len(kv_pairs) >= 3 and (any(k in nt for k in ("ficha", "datos del curso", "datos generales", "informacion general")) or len(kv_pairs) == len(non_empty)):
            _add_keyvalue_slide(sd_title, kv_pairs[:8], sd_notes)
            continue

        numbered = [l for l in sd_content if re.match(r'^\d+[\.\)]\s', l.strip())]
        dash = [l for l in sd_content if l.strip().startswith('-') or l.strip().startswith('*')]
        if len(numbered) >= 2 and len(numbered) >= len(dash) and any(k in nt for k in ("practica", "ejercicio", "actividad", "dinamica", "paso")):
            items = [_sanitize_text(re.sub(r'^\d+[\.\)]\s*', '', l.strip())) for l in numbered][:7]
            _add_numbered_slide(sd_title, items, sd_notes)
            continue

        for chunk in _split_overflow(sd):
            _add_content_slide(chunk.get("title", ""), chunk.get("content", []), chunk.get("notes", ""))

    if not has_closing:
        _add_closing_slide()

    _add_brand_outro()

    _last_pexels_credits = list(_pexels_session["credits"])

    try:
        curso_safe = re.sub(r'[^\w\s-]', '', course_name).strip().replace(' ', '_')[:30]
        output_name = f"Presentacion_{curso_safe}.pptx" if curso_safe else "Presentacion_del_Curso.pptx"
        output_path = os.path.join(GENERATED_DIR, output_name)
        prs.save(output_path)
        return output_path
    finally:
        _cleanup_temp_images()

def list_generated_docs():
    if not os.path.isdir(GENERATED_DIR):
        return []
    return sorted(os.listdir(GENERATED_DIR))
